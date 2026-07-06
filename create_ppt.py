"""
生成简洁版业务演示PPT — 减少文字、增加留白、视觉清晰
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ==================== 配色方案 ====================
C = {
    'blue':     RGBColor(0x1F, 0x77, 0xB4),
    'blue_dark':RGBColor(0x15, 0x55, 0x85),
    'blue_light':RGBColor(0xD6, 0xEB, 0xF5),
    'white':    RGBColor(0xFF, 0xFF, 0xFF),
    'black':    RGBColor(0x26, 0x27, 0x30),
    'gray':     RGBColor(0x88, 0x88, 0x88),
    'gray_bg':  RGBColor(0xF5, 0xF6, 0xF8),
    'red':      RGBColor(0xE7, 0x4C, 0x3C),
    'green':    RGBColor(0x27, 0xAE, 0x60),
    'orange':   RGBColor(0xF3, 0x9C, 0x12),
    'yellow_bg':RGBColor(0xFF, 0xF8, 0xE1),
    'green_bg': RGBColor(0xE8, 0xF8, 0xEF),
    'red_bg':   RGBColor(0xFD, 0xED, 0xEC),
    'dark_bg':  RGBColor(0x0D, 0x3B, 0x5C),
}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ==================== 工具函数 ====================

def bg(slide, color=C['white']):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def rect(slide, l, t, w, h, fill, border=None, radius=None):
    """通用矩形"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                               Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def text_box(slide, l, t, w, h, txt, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, font='微软雅黑'):
    """单行/简单文本框"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color or C['black']
    p.alignment = align; p.font.name = font
    return tf

def multi_text(slide, l, t, w, h, items):
    """多行文本 items=[(text, size, bold, color), ...]"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt, sz, bd, cl = item[0], item[1], item[2], item[3]
        p.text = txt; p.font.size = Pt(sz); p.font.bold = bd
        p.font.color.rgb = cl; p.font.name = '微软雅黑'
        p.space_after = Pt(4)
    return tf

def label(slide, l, t, w, h, txt, fill=C['blue'], text_color=C['white'], size=16):
    """圆角标签"""
    s = rect(slide, l, t, w, h, fill, radius=0.08)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(size); p.font.bold = True
    p.font.color.rgb = text_color; p.alignment = PP_ALIGN.CENTER
    p.font.name = '微软雅黑'
    return s

def circle_num(slide, l, t, d, n, fill=C['blue']):
    """带编号圆形"""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    tf = s.text_frame; p = tf.paragraphs[0]; p.text = str(n)
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = C['white']
    p.alignment = PP_ALIGN.CENTER; p.font.name = '微软雅黑'
    return s

def arrow_right(slide, l, t, w, h, color=C['gray']):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def page_num(slide, n, total=16):
    text_box(slide, 12.2, 7.1, 0.9, 0.3, f'{n}/{total}', size=10, color=C['gray'], align=PP_ALIGN.RIGHT)

def title_slide(slide, title, subtitle=''):
    """标准标题栏 + 白色背景"""
    bg(slide)
    rect(slide, 0, 0, 13.333, 1.05, C['blue'])
    # 左侧色条
    rect(slide, 0, 1.05, 0.08, 0.05, C['blue_dark'])  # tiny accent
    text_box(slide, 0.7, 0.12, 12, 0.65, title, size=30, bold=True, color=C['white'])
    if subtitle:
        text_box(slide, 0.7, 0.68, 12, 0.35, subtitle, size=13, color=RGBColor(0xBB, 0xD5, 0xEE))

def footer_bar(slide):
    """底部装饰线"""
    rect(slide, 0, 7.3, 13.333, 0.04, C['blue'])

def card(slide, l, t, w, h, title, lines, title_color=C['blue'], bg_color=C['white'], border_color=None):
    """卡片组件"""
    rect(slide, l, t, w, h, bg_color, border=border_color or RGBColor(0xE0, 0xE0, 0xE0), radius=0.12)
    text_box(slide, l + 0.25, t + 0.15, w - 0.5, 0.4, title, size=17, bold=True, color=title_color)
    multi_text(slide, l + 0.25, t + 0.55, w - 0.5, h - 0.65, lines)

def big_number(slide, l, t, number, unit, color=C['blue']):
    """大数字"""
    text_box(slide, l, t, 3, 0.8, str(number), size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
    text_box(slide, l, t + 0.7, 3, 0.35, unit, size=14, color=C['gray'], align=PP_ALIGN.CENTER)

TOTAL = 16

# ============================================================
# 第1页 — 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['blue'])
# 大圆装饰
s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
text_box(s, 1.5, 2.2, 10.3, 1.0, '注塑机台最优工艺卡分析系统',
         size=42, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
text_box(s, 2.2, 3.3, 8.9, 0.7,
         '从历史数据中自动发掘最优工艺参数\n让每一台机器都有"最优配方"',
         size=19, color=RGBColor(0xBB, 0xD5, 0xEE), align=PP_ALIGN.CENTER)
rect(s, 4.5, 4.25, 4.3, 0.025, C['white'])
text_box(s, 3.5, 4.6, 6.3, 0.5, '业务对接演示  ·  2026年7月',
         size=15, color=RGBColor(0x99, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
page_num(s, 1)

# ============================================================
# 第2页 — 一句话讲清楚
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '💡 这个系统做什么？')

label(s, 1.5, 2.2, 10.3, 1.2,
      '从历史生产数据中，自动找出每台机器跑得最好的那段时间，\n把当时的工艺参数提取出来，作为这台机器的"最优配方"推荐给你。',
      fill=C['white'], text_color=C['blue'], size=22)
# 虚线边框
s2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2))
s2.fill.background(); s2.line.color.rgb = C['blue']; s2.line.width = Pt(2)

text_box(s, 1.8, 4.0, 9.7, 1.0,
         '🎯  就像看行车记录仪 —— 找到你开车最稳、最省油的那段路，\n告诉你当时的油门、转速。我们做的就是对注塑机的"行车记录仪回放"。',
         size=18, color=C['gray'], align=PP_ALIGN.CENTER)
page_num(s, 2)

# ============================================================
# 第3页 — 痛点一览
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '🎯 解决了什么痛点？')

pains = [
    ('新师傅不熟悉机台', '→ 系统直接给出推荐值+范围'),
    ('老师傅经验无法传承', '→ 数据固化成可传承的工艺卡'),
    ('几十台机不知道哪台差', '→ 一键OEE排名，好坏立现'),
    ('OEE低但不知道原因', '→ 离线/待机/报警分类归因'),
    ('不知哪个参数影响效率', '→ 相关性量化分析'),
    ('手工翻Excel费时费力', '→ 全自动扫描，几秒出结果'),
]
for i, (pain, fix) in enumerate(pains):
    y = 1.5 + i * 0.92
    rect(s, 1.0, y, 5.0, 0.65, C['red_bg'], border=C['red'], radius=0.08)
    text_box(s, 1.2, y + 0.08, 4.6, 0.5, f'😰  {pain}', size=15, bold=True, color=C['red'])
    arrow_right(s, 6.2, y + 0.12, 0.55, 0.35, C['green'])
    rect(s, 6.9, y, 5.2, 0.65, C['green_bg'], border=C['green'], radius=0.08)
    text_box(s, 7.1, y + 0.08, 4.8, 0.5, f'✅  {fix}', size=15, bold=True, color=C['green'])

page_num(s, 3)

# ============================================================
# 第4页 — 四大模块总览
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '🗺️ 系统四大模块')

mods = [
    ('① 数据导入', '喂数据', '上传两个Excel\n自动清洗关联'),
    ('② 工艺趋势', '看参数怎么变', '趋势图 · 多机台对比\n参数-OEE相关性'),
    ('③ 稼动率分析', '看效率怎么样', '热力图 · 不稼动归因\n开合模vs稼动率'),
    ('④ 最优工艺卡', '生成最优配方 ⭐', '滑动扫描 → 最优时段\n→ 推荐参数 + 范围'),
]
for i, (name, tag, desc) in enumerate(mods):
    x = 0.5 + i * 3.2
    fill = C['red'] if i == 3 else C['blue']
    card(s, x, 1.6, 2.85, 2.6, name,
         [('', 8, False, C['black']),
          (desc, 15, False, C['black'])],
         title_color=fill, bg_color=C['white'], border_color=fill)

# 底部流程
text_box(s, 1.0, 4.8, 11.3, 0.4, '操作流程', size=18, bold=True, color=C['black'])
steps_text = ['上传两报表', '浏览趋势', '分析效率', '一键生成最优工艺卡']
for i, st in enumerate(steps_text):
    x = 0.8 + i * 3.2
    clr = C['red'] if i == 3 else C['blue']
    label(s, x, 5.4, 2.8, 0.65, st, fill=clr, size=15)
    if i < 3:
        arrow_right(s, x + 2.82, 5.52, 0.35, 0.35, C['gray'])

footer_bar(s)
page_num(s, 4)

# ============================================================
# 第5页 — 数据导入
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '📥 ① 数据导入 — 喂数据')

card(s, 0.8, 1.6, 5.5, 2.8, '📄 工艺趋势报表', [
    ('每台机台 · 每小时的工艺参数值', 15, False, C['black']),
    ('', 6, False, C['black']),
    ('射胶参数  压力 / 速度 / 位置 × 5段', 13, False, C['gray']),
    ('保压参数  压力 / 速度 / 时间 × 3段', 13, False, C['gray']),
    ('温度参数  一段 ~ 六段', 13, False, C['gray']),
    ('其他      储料背压、成型周期等', 13, False, C['gray']),
])

card(s, 7.0, 1.6, 5.5, 2.8, '📄 稼动率小时报表', [
    ('每台机台 · 每小时的运行效率', 15, False, C['black']),
    ('', 6, False, C['black']),
    ('设备稼动率 (OEE)', 13, False, C['gray']),
    ('有效生产时长', 13, False, C['gray']),
    ('离线 / 待机 / 报警时长（停机归因）', 13, False, C['gray']),
    ('开合模次数', 13, False, C['gray']),
])

items = [
    '日期格式自动标准化',
    '两表自动关联匹配（机台号+日期+时间段）',
    '异常数据自动标记（负数 / 超高值）',
    '数据自动缓存，下次打开无需重新上传',
]
for i, item in enumerate(items):
    label(s, 0.8 + (i % 2) * 6.2, 4.9 + (i // 2) * 0.7, 5.8, 0.55, f'✔  {item}',
         fill=C['blue_light'], text_color=C['blue_dark'], size=13)

footer_bar(s)
page_num(s, 5)

# ============================================================
# 第6页 — 工艺趋势
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '📈 ② 工艺趋势 — 看参数怎么变')

modes = [
    ('📈 单机台趋势', '选一台机器\n看参数随时间变化曲线'),
    ('🔄 多机台对比', '选多台机器\n对比同一个参数差异'),
    ('🔗 参数-OEE相关性', '量化每个参数\n对效率的影响大小'),
]
for i, (name, desc) in enumerate(modes):
    x = 0.8 + i * 4.2
    card(s, x, 1.6, 3.7, 1.8, name, [(desc, 16, False, C['black'])],
         bg_color=C['white'], border_color=C['blue'])

text_box(s, 0.8, 3.8, 11.5, 0.5, '💼 业务价值', size=20, bold=True, color=C['blue'])
vals = ['波动大的参数 → 重点管控', '同型号差异大 → 排查原因', '强相关参数 → 调机优先关注']
for i, v in enumerate(vals):
    label(s, 0.8 + i * 4.2, 4.5, 3.7, 0.65, v, fill=C['blue_light'], text_color=C['blue_dark'], size=15)

# 解释相关性
text_box(s, 1.0, 5.6, 11.3, 1.2,
         '🔗 相关性解读：绿色=正相关（参数大→OEE高）  |  红色=负相关（参数大→OEE低）\n'
         '    绝对值越大 → 影响越大。比如"成型周期"与OEE负相关：周期越长，效率越低。',
         size=14, color=C['gray'])
footer_bar(s)
page_num(s, 6)

# ============================================================
# 第7页 — 稼动率分析
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '⏱️ ③ 稼动率分析 — 看效率怎么样')

charts = [
    ('🔥 稼动率热力图', '颜色 = OEE高低\n红=差 黄=中 绿=好\n一眼定位问题时段'),
    ('🥧 不稼动归因饼图', '离线 待机 报警\n三类停机原因占比\n搞清楚OEE损失的第一责任人'),
    ('📊 按天对比柱状图', '每天每小时一组柱子\n不同日期颜色区分\n发现黄金时段和低谷'),
    ('🔵 开合模 vs 稼动率', '产量 vs 效率的关系\n蓝点=正常 红X=异常\n排查传感器故障'),
]
for i, (title, desc) in enumerate(charts):
    x = 0.4 + i * 3.2
    card(s, x, 1.6, 3.0, 2.6, title, [(desc, 14, False, C['black'])],
         bg_color=C['white'], border_color=C['blue'])

text_box(s, 0.9, 4.7, 11.5, 0.8,
         '💡 核心价值："知道OEE低" → "知道为什么低、什么时间最低、哪台机器最差"',
         size=20, bold=True, color=C['blue'])
footer_bar(s)
page_num(s, 7)

# ============================================================
# 第8页 — 最优工艺卡总览 ⭐
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, 13.333, 1.05, C['red'])
text_box(s, 0.7, 0.12, 12, 0.65, '🏆 ④ 最优工艺卡 — 核心功能',
         size=30, bold=True, color=C['white'])

# 三大步骤
text_box(s, 1.0, 1.5, 11.3, 0.5, '三步生成最优工艺卡', size=22, bold=True, color=C['black'])
steps = [
    ('① 扫描', '逐小时扫描历史数据\n找出所有连续OEE达标的时段', '🔍'),
    ('② 选取', '挑出持续时间最长的窗口\n（同样长选OEE最高的）', '✅'),
    ('③ 提取', '取窗口内参数的中位数\n作为推荐值 + 建议范围', '📋'),
]
for i, (title, desc, icon) in enumerate(steps):
    x = 1.0 + i * 4.2
    label(s, x, 2.3, 3.6, 1.6, '', fill=C['white'], text_color=C['black'])
    step_color = C['blue'] if i < 2 else C['green']
    text_box(s, x + 0.3, 2.4, 3.0, 0.6, f"{icon}  {title}", size=20, bold=True, color=step_color)
    text_box(s, x + 0.3, 2.9, 3.0, 0.9, desc, size=14, color=C['gray'])

# 输出示例
text_box(s, 1.0, 4.3, 11.3, 0.5, '输出示例：一张完整的最优工艺卡', size=18, bold=True, color=C['black'])

# 简化的示例表
example_data = [
    ['参数名', '推荐值⭐', '建议下限', '建议上限'],
    ['射胶压力一段', '85.5 bar', '84.0', '87.3'],
    ['保压压力一段', '65.0 bar', '63.0', '67.0'],
    ['温度二段值', '235.0 ℃', '233.0', '238.0'],
    ['成型周期（s）', '28.5', '27.8', '29.3'],
]
tbl_shape = s.shapes.add_table(len(example_data), 4, Inches(1.0), Inches(4.9), Inches(11.3), Inches(2.0))
tbl = tbl_shape.table
for ri, row in enumerate(example_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = C['blue'] if ri == 0 else (C['gray_bg'] if ri % 2 == 0 else C['white'])
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15 if ri == 0 else 14)
            p.font.bold = (ri == 0)
            p.font.color.rgb = C['white'] if ri == 0 else C['black']
            p.alignment = PP_ALIGN.CENTER
            p.font.name = '微软雅黑'
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.5)

footer_bar(s)
page_num(s, 8)

# ============================================================
# 第9页 — 算法核心逻辑
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '🔍 核心算法：怎么找到跑得最好的那段？')

text_box(s, 0.8, 1.4, 11.5, 0.4, '条件：连续 ≥ 12 小时，每小时 OEE ≥ 90%', size=17, bold=True, color=C['black'])

# 场景1
rect(s, 0.8, 2.1, 11.5, 0.9, C['red_bg'], border=C['red'], radius=0.08)
text_box(s, 1.1, 2.2, 11, 0.3, '场景①：找到了连续 6 小时达标 → 不够 12 小时 ❌ 放弃', size=16, bold=True, color=C['red'])
text_box(s, 1.1, 2.55, 11, 0.4, '92%  95%  91%  93%  94%  90%  78%  62%  ...     只有6小时达标，不满足条件', size=15, color=C['gray'])

# 场景2
rect(s, 0.8, 3.3, 11.5, 0.9, C['green_bg'], border=C['green'], radius=0.08)
text_box(s, 1.1, 3.4, 11, 0.3, '场景②：找到了连续 13 小时达标 → ≥ 12 小时 ✅ 这就是最优窗口！', size=16, bold=True, color=C['green'])
text_box(s, 1.1, 3.75, 11, 0.4, '92%  95%  91%  93%  94%  90%  91%  93%  95%  92%  91%  90%  93%   连续13小时！', size=15, color=C['gray'])

# 三个关键设计
text_box(s, 0.8, 4.7, 11.5, 0.4, '三个关键设计保障结果可靠', size=20, bold=True, color=C['black'])

designs = [
    ('数据缺口自动检测', '中间缺失超过1.5小时会自动断开，确保"连续"是真连续'),
    ('异常OEE裁剪处理', '负数或超高值(>150%)先裁剪到0~100%，避免传感器故障干扰判定'),
    ('多窗口智能选优', '多窗口时优先选持续时间最长的；同样长选平均OEE最高的'),
]
for i, (d_title, d_desc) in enumerate(designs):
    circle_num(s, 1.0, 5.4 + i * 0.65, 0.35, i + 1, C['blue'])
    text_box(s, 1.55, 5.35 + i * 0.65, 10, 0.55, f'{d_title}：{d_desc}', size=14, color=C['black'])

footer_bar(s)
page_num(s, 9)

# ============================================================
# 第10页 — 为什么用中位数
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '📊 为什么推荐值用"中位数"而不是"平均值"？')

text_box(s, 0.8, 1.4, 11.5, 0.4,
         '场景：窗口内某参数13个读数，其中1个是传感器瞬间跳变（异常高值）',
         size=15, color=C['gray'])

label(s, 1.5, 2.0, 10.3, 0.7,
      '82  ·  84  ·  85  ·  85  ·  85  ·  86  ·  86  ·  87  ·  87  ·  88  ·  89  ·  90  ·  120(异常↑)',
      fill=C['blue_light'], text_color=C['blue_dark'], size=16)

# 对比
card(s, 1.0, 3.2, 5.3, 2.8, '❌ 用平均值',
     [('总和 / 13 = 88.8', 22, True, C['red']),
      ('', 8, False, C['black']),
      ('被 120 拉高了近 3 个 bar！', 16, True, C['red']),
      ('按 88.8 去设 → 压力偏大', 14, False, C['gray']),
      ('可能导致飞边、浪费材料', 14, False, C['gray'])],
     title_color=C['red'], border_color=C['red'])

card(s, 6.9, 3.2, 5.3, 2.8, '✅ 用中位数',
     [('排序取中间 = 86.0', 22, True, C['green']),
      ('', 8, False, C['black']),
      ('完全不受 120 的影响！', 16, True, C['green']),
      ('按 86.0 去设 → 真实典型值', 14, False, C['gray']),
      ('安全、稳定、靠谱', 14, False, C['gray'])],
     title_color=C['green'], border_color=C['green'])

label(s, 2.5, 6.5, 8.3, 0.6,
      '💡 中位数 = 不受极端值干扰 = 代表"大多数时候"的真实值',
      fill=C['blue'], text_color=C['white'], size=18)
footer_bar(s)
page_num(s, 10)

# ============================================================
# 第11页 — 演示路径
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '🎬 15分钟演示路径')

demo = [
    ('01', '系统主页', '展示OEE排名', C['blue']),
    ('02', '数据导入', '演示上传+自动处理', C['blue']),
    ('03', '工艺趋势', '参数曲线+相关性', C['blue']),
    ('04', '稼动率分析', '热力图+归因饼图', C['blue']),
    ('05', '最优工艺卡', '⭐ 设条件→分析→看结果', C['red']),
    ('06', '总结', '场景+价值回顾', C['blue']),
]
for i, (num, module, action, clr) in enumerate(demo):
    x = 0.6 + i * 2.15
    label(s, x, 1.6, 1.9, 1.5, '', fill=C['white'], text_color=C['black'])
    text_box(s, x + 0.15, 1.7, 1.6, 0.5, num, size=28, bold=True, color=clr, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.15, 2.2, 1.6, 0.4, module, size=15, bold=True, color=clr, align=PP_ALIGN.CENTER)
    text_box(s, x + 0.15, 2.55, 1.6, 0.45, action, size=12, color=C['gray'], align=PP_ALIGN.CENTER)
    if i < 5:
        arrow_right(s, x + 1.9, 2.2, 0.22, 0.3, C['gray'])

# 演示要点
text_box(s, 0.8, 3.6, 11.5, 0.5, '演示重点', size=20, bold=True, color=C['black'])
tips = [
    '第1步：OEE排名让大家看到"原来差异这么大"',
    '第3步：相关性分析让大家看到"原来这个参数这么重要"',
    '第5步：全选所有机台→一键分析→逐个查看→下载导出（一气呵成）',
    '最后：强调"以后换模调参，先查系统"的习惯养成',
]
for i, tip in enumerate(tips):
    text_box(s, 1.2, 4.3 + i * 0.55, 10.5, 0.5, f'• {tip}', size=15, color=C['black'])

footer_bar(s)
page_num(s, 11)

# ============================================================
# 第12页 — 应用场景
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '💼 业务应用场景')

scenes = [
    ('🆕', '新品上线\n设定首版参数', '历史最优值为基线\n建议范围内微调\n减少试错成本'),
    ('🔧', '效率突然变差\n需要排查原因', '定位波动大的参数\n找出强相关参数\n对比当前值 vs 推荐值'),
    ('📊', '多台同型号\n效率差异大', 'OEE排名看谁好谁差\n参数对比找差异点\n最优工艺卡横向对比'),
    ('📝', '周会月会\n汇报材料', 'OEE排名图 → 进PPT\n不稼动饼图 → 分析原因\n工艺卡 → 展示改善成果'),
    ('👨‍🏫', '新员工培训\n快速上手', '最优工艺卡 = 标准参考\n减少"跟师傅调3个月"\n降低人员流失冲击'),
]
for i, (icon, scene, detail) in enumerate(scenes):
    x = 0.3 + i * 2.6
    card(s, x, 1.6, 2.35, 4.0, f'{icon}  {scene}',
         [(detail, 14, False, C['black'])],
         bg_color=C['white'], border_color=C['blue'])

footer_bar(s)
page_num(s, 12)

# ============================================================
# 第13页 — 业务对齐
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '🤝 业务对齐：需要达成共识的关键认知')

aligns = [
    ('推荐值 ≠ 唯一正确答案', C['blue'],
     '推荐值是"历史最好时段的常用值"，不是绝对的。不同产品/模具/原料需微调，把它当"起点"不是"终点"。'),
    ('数据越多 → 推荐越准', C['green'],
     '3天数据能找到一个"看起来不错"的窗口；7天数据"比较可靠"；30天数据"非常靠谱"。建议持续积累。'),
    ('相关性 ≠ 因果性', C['orange'],
     '系统告诉你"参数X和OEE强相关"，不等于"调X一定改善OEE"。需结合工艺知识判断。它是"帮你发现候选参数"的工具。'),
    ('数据质量决定结果质量', C['red'],
     '数采系统经常断、经常出异常 → 分析结果可信度下降。系统已做了异常标记，但源头准确性最关键。'),
]
for i, (title, clr, desc) in enumerate(aligns):
    y = 1.5 + i * 1.35
    circle_num(s, 0.8, y + 0.1, 0.4, i + 1, clr)
    text_box(s, 1.4, y, 3.2, 0.5, title, size=18, bold=True, color=clr)
    text_box(s, 1.4, y + 0.5, 11, 0.7, desc, size=14, color=C['black'])

footer_bar(s)
page_num(s, 13)

# ============================================================
# 第14页 — FAQ
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '❓ 常见问题 FAQ')

faqs = [
    ('Q1  系统能替代老师傅的经验吗？',
     '不能，但能辅助。系统把数据里的"最优值"找出来，师傅做判断和微调。系统+经验=最优组合。'),
    ('Q2  半夜无人值守的数据更可靠吗？',
     '恰恰更可靠。无人时段调试少、干扰少，参数更稳定。高OEE说明参数组合本身就好。'),
    ('Q3  换模具/换材料后还能用吗？',
     '不能直接用，但可作为参考基线。积累新数据后重新分析，生成对应版本的最优工艺卡。'),
    ('Q4  有些机台找不到最优窗口怎么办？',
     '说明这台机在数据期内没有真正稳定运行过。降低标准（80%+6h）再试；仍找不到→需要重点整治。'),
    ('Q5  多久更新一次？',
     '建议每周积累新数据后跑一次。重大变更（换模/大修/换料）后也重新跑。全车间几十台也就几十秒。'),
]
for i, (q, a) in enumerate(faqs):
    y = 1.45 + i * 1.1
    text_box(s, 0.8, y, 11.5, 0.35, q, size=17, bold=True, color=C['blue'])
    text_box(s, 1.1, y + 0.4, 11.2, 0.55, a, size=14, color=C['black'])

footer_bar(s)
page_num(s, 14)

# ============================================================
# 第15页 — 速记卡
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, '📌 关键信息速记卡')

card_data = [
    ('需要什么数据？', '两个Excel：\n工艺参数 + 稼动率'),
    ('怎么判"跑得好"？', 'OEE ≥ 90%\n连续 ≥ 12小时'),
    ('怎么选最优时段？', '找连续达标时段\n选持续时间最长'),
    ('推荐值怎么算？', '窗口内中位数\n不受极端值影响'),
    ('建议范围怎么算？', '下限 Q25\n上限 Q75'),
    ('能分析多少台？', '不限\n一键全车间扫描'),
    ('多久更新一次？', '每周一次\n重大变更后更新'),
    ('结果能导出吗？', '可下载CSV\nExcel直接打开'),
]
for i, (q, a) in enumerate(card_data):
    x = 0.5 + (i % 4) * 3.15
    y = 1.6 + (i // 4) * 2.7
    label(s, x, y, 2.9, 0.55, q, fill=C['blue'], size=14)
    label(s, x, y + 0.7, 2.9, 1.2, a, fill=C['blue_light'], text_color=C['blue_dark'], size=13)

footer_bar(s)
page_num(s, 15)

# ============================================================
# 第16页 — 结束页
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, C['blue'])
text_box(s, 1.5, 1.8, 10.3, 1.0, '🏭', size=56, color=C['white'], align=PP_ALIGN.CENTER)
text_box(s, 1.5, 2.7, 10.3, 1.0, '总结', size=38, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

label(s, 1.5, 4.0, 10.3, 1.2,
      '这个系统不替代你的经验，\n它帮你把隐藏在几十万行历史数据里的"最优跑法"找出来，\n变成一张看得见、用得上、可传承的工艺卡。',
      fill=C['blue_dark'], text_color=C['white'], size=21)

rect(s, 4.5, 5.5, 4.3, 0.025, C['white'])
text_box(s, 2.0, 5.8, 9.3, 0.7, '感谢聆听 · 欢迎提问 🙋', size=20,
         color=RGBColor(0xBB, 0xD5, 0xEE), align=PP_ALIGN.CENTER)
page_num(s, 16)

# ============================================================
# 保存
# ============================================================
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), '注塑机台最优工艺卡_业务演示_v2.pptx')
prs.save(output)
print(f'PPT saved: {output}')
print(f'Total slides: {len(prs.slides)}')